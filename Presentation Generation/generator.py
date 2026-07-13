# generator.py
from llm import call_model
from io import BytesIO
import pptx, re
from config import TITLE_FONT_SIZE, SLIDE_FONT_SIZE

def generate_slide_titles(topic, category, slide_number) -> list[str]:
    prompt = f"Generate three different title options for slide number '{slide_number}' for the topic '{topic}' with category '{category}'."
    raw = call_model(prompt)
    titles = [line.replace('**','') for line in raw.split('\n') if line.strip() and line.strip()[0].isdigit()]
    titles = [t.split('. ',1)[1] if '. ' in t else t for t in titles]
    titles = [t.split(': ')[1] if ': ' in t else t for t in titles]
    return [t.replace('"','').strip() for t in titles]

def generate_slide_content(topic, title, category, description="") -> tuple[str, str]:
    prompt = f"Generate 2 versions of content for slide '{title}', topic '{topic}', category '{category}', description '{description}'. Give 3 lines of detail and 3-4 bullet points. No Image Suggestions, Color Scheme, Fonts, or Slide Title."
    raw = call_model(prompt)
    clean = raw.replace('**','').replace('*','').replace('• ','')
    parts = clean.replace('Version 1:','').split('Version 2:')
    v1 = parts[0].strip() if len(parts) > 0 else clean
    v2 = parts[1].strip() if len(parts) > 1 else clean
    return v1, v2

def create_presentation(topic, slide_titles, slide_contents) -> BytesIO:
    prs = pptx.Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for title, content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para.font.size = SLIDE_FONT_SIZE

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer