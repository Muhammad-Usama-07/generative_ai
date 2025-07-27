import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
from ctransformers import AutoModelForCausalLM
from dotenv import load_dotenv
import pandas as pd
import numpy as np
from textwrap import dedent
from typing import List
from fastembed import TextEmbedding
from groq import Groq
from groq.types.chat import ChatCompletionMessage
import streamlit as st
import re
load_dotenv()


# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)



# local_llm = r'llama-2-7b-chat.Q4_K_M.gguf'

# llm = AutoModelForCausalLM.from_pretrained(local_llm, model_type='llama')


# Set up the Groq API key and client
os.environ["GROQ_API_KEY"] = ""
client = Groq()

MODEL = "llama-3.1-70b-versatile"
# MODEL = "7b-chat-q4_K_M"

TEMPERATURE = 0
MAX_TOKENS = 4096

def call_model(query: str, messages: List) -> ChatCompletionMessage:
    messages.append(
        {
            "role": "user",
            "content": query,
        },
    )
    response = client.chat.completions.create(
        model=MODEL,
        messages=messages,
        temperature=TEMPERATURE,
        max_tokens=MAX_TOKENS,
    )
    message = response.choices[0].message
    messages.append(message)
    return message


SYSTEM_PROMPT = '''You are an AI-Powered Presentation Generator, designed to generate 2 best detailed version of professional PowerPoint presentations with complete explanation of the topics and the topic that user provide must divide in slide by slide sections.
Your primary task is to generate content based on user inputs related to presentation slides. You're always helpful to generate  presenetation only based on the provided topic. If you don't unserstand the topic - just reply with an excuse that you
don't know. Keep your content brief and detailed. Be kind and respectful.'''
# SYSTEM_PROMPT = ''


def generate_slide_titles(topic, no_of_slides, presentation_category, slide_number):
    # prompt = f"Generate {no_of_slides} slide title for the topic '{topic}' with the category {presentation_category}."
    prompt = f"Generate three different types of title for slide number '{slide_number}' for the topic '{topic}' with the category {presentation_category}."
    
    # prompt = f"Generate heading for slide one of the presentation for the topic '{topic}' with the category {presentation_category}, and take heading idea from this text '{step2_heading}'."
    # prompt = f"i want to make presentation of {no_of_slides} number of slides, now generate me title of slide one for the topic '{topic}' with the category {presentation_category}, and take heading idea from this text '{step2_heading}"



    # response = llm(prompt)


    if prompt:
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
        ]

        # with st.spinner("Generating response..."):
        bot_response = call_model(prompt, messages)

        # st.write("**Bot Response:**")
        # st.write(bot_response.content)

    
    return bot_response.content.split("\n")

def generate_slide_content(topic, title, presentation_category, step2_description):
    # prompt = f"Generate content for the slide: '{slide_title}' with the category {presentation_category}. dont give 'Image Suggestions', 'Color Scheme', 'Fonts', and 'Slide Title'"
    # prompt = f"Generate content for the slide: '{title}' with the category {presentation_category}, description {step2_description}, and topic {topic}. dont give 'Image Suggestions', 'Color Scheme', 'Fonts', 'Key Points to Discuss',  and 'Slide Title'"
    prompt = f"Generate 2 best versions of the content for the slide: '{title}' with the category {presentation_category}, description {step2_description}, and topic {topic}. remember to just only give 3 lines related to the detail and 3-4 bullet points and dont give 'Image Suggestions', 'Color Scheme', 'Fonts', 'Key Points to Discuss',  and 'Slide Title just give detail and keypoints"

    
    # print(prompt)
    # response = llm(prompt)
    # print(f'\nresponse of the title {slide_title}\n', response)
    if prompt:
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
        ]

        # with st.spinner("Generating response..."):
        bot_response = call_model(prompt, messages)
    # print('\n\nbot_response:\n')
    # print(bot_response.content)

    return bot_response.content.replace('**', '').replace('*', '').replace('• ', '').replace('Detail:', '').replace('Detail: \n', '').replace('Key Points:\n', '')


def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    # prs.save(f"generated_ppt/{topic}_presentation.pptx")
    prs.save(f"test_generated_presentation.pptx")


from pprint import pprint 
topic = input("Enter the topic for your presentation:")
topic = re.sub(r'[^a-zA-Z\s]', '', topic)
category_list  = presentation_types = [
    '1. Business presentations',
    '2. Motivational presentation',
    '3. Persuasive presentations',
    '4. Demonstrative presentations',
    '5. Informative',
    '6. Storytelling',
    '7. Academic presentations',
    '8. Sales presentation',
    '9. Instructional',
    '10. Progress presentations',
    '11. Coach style',
    '12. Connector style',
    '13. Decision-making presentations',
    '14. Instructor style',
    '15. Problem solving'
]

print('\n\nselect category: ', category_list)
presentation_category_index = int(input("\nEnter a number between 1 and 15: :"))

pattern = r'[0-9]'
selected_category = re.sub(pattern, '', category_list[presentation_category_index]).replace('. ','')
print('*** seleted presentation_category: ',selected_category )
no_of_slides = int(input("No. of slides:"))
# presen_details = {}
presen_titles = []
presen_content = []

for i in range(1, no_of_slides+1):
    slide_number = i
    if i == no_of_slides:
        # slide_number = 'second last'
        suggested_titles = generate_slide_titles(topic, no_of_slides, selected_category, 'second last')
    else:
        suggested_titles = generate_slide_titles(topic, no_of_slides, selected_category, slide_number)
    filtered_slide_titles = [item.replace('**', '') for item in suggested_titles if item.strip() != '']
    filtered_slide_titles = [title for title in filtered_slide_titles if title[0].isdigit()]
    filtered_slide_titles = [title.split('. ', 1)[1] for title in filtered_slide_titles if title[0].isdigit()]
    filtered_slide_titles = [title.split(': ')[1] if len(title.split(': ')) > 1 else title for title in filtered_slide_titles ]
    filtered_slide_titles = [item.replace('"', '') for item in filtered_slide_titles]
    # print('\nsuggested_titles: ', filtered_slide_titles)
    print(f'\n\nPick title for slide {slide_number}: ', filtered_slide_titles)
    filtered_slide_titles_index = int(input("\nEnter your choice:"))
    selected_slide_title = filtered_slide_titles[filtered_slide_titles_index]
    des = input("\nEnter some description or skip:")
    print('selected_slide_title: ', selected_slide_title)
    text_versions = generate_slide_content(topic, selected_slide_title, selected_category, des).replace('Version 1:', '').split('Version 2:')
    print('\nSelect content verisons: ', text_versions)
    text_versions_index = int(input("\nEnter with 0 or 1:"))
    select_version = text_versions[text_versions_index]
    print('\n\nSelected veriosn: ', select_version)
    presen_titles.append(selected_slide_title)
    presen_content.append(select_version)


    # presen_details[f'slide_{slide_number}'] = {'titles': filtered_slide_titles}

# print('\npresen_details: ', presen_details)

create_presentation(topic, presen_titles, presen_content)
print("Presentation generated successfully!")







# # slide_content = slide_content
# print('\n\nslide_content V1: ', v1)
# print('\n\nslide_content V2 Version V2: ', v2)

# # create_presentation(topic, slide_title, slide_content)
# print('------- presentation created ------------')
# # step1_button = st.button("Go to slide details")

